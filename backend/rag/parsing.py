"""Local document parsing.

Extracts text and structure from the formats an industrial site actually
holds: PDF, DOCX, XLSX/CSV, PPTX, plain text. Everything runs on-host with no
service calls. Pages/sheets/slides are preserved as *locations* so retrieved
evidence can cite "page 4" or "sheet Inspection!A1:F20" rather than an opaque
chunk id.
"""

from __future__ import annotations

import csv
import io
import re
from dataclasses import dataclass
from pathlib import Path


@dataclass
class ParsedSegment:
    """One addressable region of a source document."""

    text: str
    location: str


@dataclass
class ParsedDocument:
    segments: list[ParsedSegment]
    media_type: str
    page_count: int
    parser: str
    warnings: list[str]

    @property
    def full_text(self) -> str:
        return "\n\n".join(segment.text for segment in self.segments if segment.text.strip())


class ParsingError(RuntimeError):
    """Raised when a document cannot be parsed locally."""


HEADING_PATTERN = re.compile(r"^(#{1,6})\s+(.*\S)\s*$")


def _parse_text(path: Path) -> ParsedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")

    # Markdown carries its own structure; splitting on headings lets retrieved
    # evidence cite "section: 4. Corrosion Rate" instead of an opaque offset.
    if path.suffix.lower() == ".md":
        segments: list[ParsedSegment] = []
        current_heading = "preamble"
        buffer: list[str] = []

        def flush() -> None:
            body = "\n".join(buffer).strip()
            if body:
                segments.append(
                    ParsedSegment(text=body, location=f"section: {current_heading}")
                )

        for line in text.splitlines():
            match = HEADING_PATTERN.match(line)
            if match:
                flush()
                buffer = []
                current_heading = match.group(2).strip()
                buffer.append(line)
            else:
                buffer.append(line)
        flush()

        if segments:
            return ParsedDocument(
                segments=segments,
                media_type="text/markdown",
                page_count=len(segments),
                parser="markdown",
                warnings=[],
            )

    return ParsedDocument(
        segments=[ParsedSegment(text=text, location="whole document")],
        media_type="text/plain",
        page_count=1,
        parser="plaintext",
        warnings=[],
    )


def extract_title(path: Path) -> str | None:
    """Best-effort document title from its own content, not its filename."""
    if path.suffix.lower() not in {".md", ".txt"}:
        return None
    try:
        for line in path.read_text(encoding="utf-8", errors="replace").splitlines()[:40]:
            match = HEADING_PATTERN.match(line)
            if match and len(match.group(1)) == 1:
                return match.group(2).strip()
            stripped = line.strip()
            if stripped and not stripped.startswith(("#", "*", "-", "|", ">")):
                return stripped[:120]
    except OSError:
        return None
    return None


def _parse_pdf(path: Path) -> ParsedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover
        raise ParsingError("pypdf is not installed") from exc

    reader = PdfReader(str(path))
    segments: list[ParsedSegment] = []
    warnings: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # damaged page, keep going
            text = ""
            warnings.append(f"page {index}: extraction failed ({exc})")
        if text.strip():
            segments.append(ParsedSegment(text=text, location=f"page {index}"))

    if not segments:
        warnings.append(
            "No embedded text found. This PDF is likely a scan; route it to the "
            "vision model for OCR instead of the text parser."
        )
    return ParsedDocument(
        segments=segments,
        media_type="application/pdf",
        page_count=len(reader.pages),
        parser="pypdf",
        warnings=warnings,
    )


def _parse_docx(path: Path) -> ParsedDocument:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover
        raise ParsingError("python-docx is not installed") from exc

    document = docx.Document(str(path))
    segments: list[ParsedSegment] = []
    buffer: list[str] = []
    section_index = 1
    current_heading = "body"

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if paragraph.style is not None and str(paragraph.style.name).startswith("Heading"):
            if buffer:
                segments.append(
                    ParsedSegment(text="\n".join(buffer), location=f"section: {current_heading}")
                )
                buffer = []
                section_index += 1
            current_heading = text
        buffer.append(text)

    if buffer:
        segments.append(ParsedSegment(text="\n".join(buffer), location=f"section: {current_heading}"))

    for table_index, table in enumerate(document.tables, start=1):
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells) for row in table.rows
        ]
        if rows:
            segments.append(
                ParsedSegment(text="\n".join(rows), location=f"table {table_index}")
            )

    return ParsedDocument(
        segments=segments,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        page_count=max(1, section_index),
        parser="python-docx",
        warnings=[],
    )


def _parse_xlsx(path: Path) -> ParsedDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise ParsingError("openpyxl is not installed") from exc

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    segments: list[ParsedSegment] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [("" if value is None else str(value)) for value in row]
            if any(value.strip() for value in values):
                rows.append(" | ".join(values))
        if rows:
            segments.append(
                ParsedSegment(text="\n".join(rows), location=f"sheet '{sheet.title}'")
            )
    workbook.close()
    return ParsedDocument(
        segments=segments,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        page_count=len(segments) or 1,
        parser="openpyxl",
        warnings=[],
    )


def _parse_csv(path: Path) -> ParsedDocument:
    raw = path.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(raw))
    rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]
    return ParsedDocument(
        segments=[ParsedSegment(text="\n".join(rows), location="whole file")] if rows else [],
        media_type="text/csv",
        page_count=1,
        parser="csv",
        warnings=[],
    )


def _parse_pptx(path: Path) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ParsingError("python-pptx is not installed") from exc

    presentation = Presentation(str(path))
    segments: list[ParsedSegment] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if lines:
            segments.append(ParsedSegment(text="\n".join(lines), location=f"slide {index}"))
    return ParsedDocument(
        segments=segments,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        page_count=len(presentation.slides),
        parser="python-pptx",
        warnings=[],
    )


PARSERS = {
    ".txt": _parse_text,
    ".md": _parse_text,
    ".log": _parse_text,
    ".json": _parse_text,
    ".yaml": _parse_text,
    ".yml": _parse_text,
    ".py": _parse_text,
    ".js": _parse_text,
    ".ts": _parse_text,
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".xls": _parse_xlsx,
    ".csv": _parse_csv,
    ".pptx": _parse_pptx,
}

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif", ".gif"}


def is_image(path: Path) -> bool:
    return path.suffix.lower() in IMAGE_SUFFIXES


def parse_document(path: Path) -> ParsedDocument:
    """Parse a local document into addressable text segments."""
    if not path.exists():
        raise ParsingError(f"File not found: {path}")
    suffix = path.suffix.lower()
    if suffix in IMAGE_SUFFIXES:
        raise ParsingError(
            f"'{path.name}' is an image. Images are understood by the vision "
            "model, not the text parser."
        )
    parser = PARSERS.get(suffix)
    if parser is None:
        raise ParsingError(f"No local parser registered for '{suffix}' files")
    return parser(path)
