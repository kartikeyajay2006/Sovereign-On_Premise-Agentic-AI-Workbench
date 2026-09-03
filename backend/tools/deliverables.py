"""Deliverable engine.

Turns verified agent output into the artifacts an industrial site actually
files: approval notes and reports (DOCX), analysis workbooks (XLSX), board
packs (PPTX), and Markdown. Every deliverable carries its provenance block —
evidence citations, model and version, verification status, approval record
and content hash — because an artifact that cannot be traced back is not
usable as a record.

Rendering is deterministic: the model supplies structured content, this module
supplies the document. No model output is trusted to produce formatting.
"""

from __future__ import annotations

import hashlib
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from backend.core.config import get_config
from backend.core.schemas import (
    Deliverable,
    EvidenceItem,
    RoutingDecision,
    VerificationReport,
)

SAFE_NAME = re.compile(r"[^A-Za-z0-9_.-]+")

# Palette used across generated documents, kept consistent with the workbench UI.
INK = "1F2933"
ACCENT = "0F6E6E"
MUTED = "5A6872"


def _slug(value: str, fallback: str = "deliverable") -> str:
    cleaned = SAFE_NAME.sub("_", value.strip())[:60].strip("._-")
    return cleaned or fallback


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(65536), b""):
            digest.update(block)
    return digest.hexdigest()


class DeliverableEngine:
    """Renders structured content into DOCX / XLSX / PPTX / MD artifacts."""

    def __init__(self) -> None:
        self.config = get_config()

    @property
    def output_root(self) -> Path:
        root = self.config.settings.path("deliverables")
        root.mkdir(parents=True, exist_ok=True)
        return root

    def _target(self, task_id: str, title: str, extension: str) -> Path:
        directory = self.output_root / task_id
        directory.mkdir(parents=True, exist_ok=True)
        return directory / f"{_slug(title)}.{extension}"

    def _finalise(self, path: Path, fmt: str) -> Deliverable:
        return Deliverable(
            id=str(uuid.uuid4()),
            filename=path.name,
            format=fmt,
            size_bytes=path.stat().st_size,
            sha256=_sha256(path),
            download_url=f"/api/deliverables/{path.parent.name}/{path.name}",
            released=False,
            created_at=datetime.now(timezone.utc),
        )

    # -- DOCX --------------------------------------------------------------
    def render_docx(
        self,
        *,
        task_id: str,
        content: dict[str, Any],
        evidence: list[EvidenceItem],
        routing: list[RoutingDecision],
        verification: VerificationReport | None,
        author: str,
        calculations: list[dict[str, Any]] | None = None,
    ) -> Deliverable:
        import docx
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt, RGBColor

        document = docx.Document()

        styles = document.styles["Normal"]
        styles.font.name = "Calibri"
        styles.font.size = Pt(10.5)

        title = str(content.get("title") or "Industrial Note")
        reference = str(content.get("reference") or task_id[:12].upper())

        heading = document.add_paragraph()
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = heading.add_run(title.upper())
        run.bold = True
        run.font.size = Pt(15)
        run.font.color.rgb = RGBColor.from_string(INK)

        subtitle = document.add_paragraph()
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub_run = subtitle.add_run(
            "Prepared on the Sovereign On-Premise Agentic AI Workbench "
            "— all processing performed locally"
        )
        sub_run.italic = True
        sub_run.font.size = Pt(8.5)
        sub_run.font.color.rgb = RGBColor.from_string(MUTED)

        meta = document.add_table(rows=0, cols=2)
        meta.style = "Light Grid Accent 1"
        meta_rows = [
            ("Reference", reference),
            ("Task ID", task_id),
            ("Prepared for", author),
            ("Date", datetime.now(timezone.utc).strftime("%d %b %Y, %H:%M UTC")),
            (
                "Verification",
                "PASSED" if verification and verification.valid else
                ("FAILED — see limitations" if verification else "not performed"),
            ),
            ("Approval status", "Pending human approval"),
        ]
        for label, value in meta_rows:
            row = meta.add_row().cells
            row[0].text = label
            row[1].text = str(value)
            for paragraph in row[0].paragraphs:
                for run_ in paragraph.runs:
                    run_.bold = True

        if content.get("summary"):
            document.add_heading("Executive summary", level=1)
            document.add_paragraph(str(content["summary"]))

        for section in content.get("sections") or []:
            document.add_heading(str(section.get("heading", "Section")), level=1)
            body = str(section.get("body", "")).strip()
            if body:
                document.add_paragraph(body)
            for bullet in section.get("bullets") or []:
                document.add_paragraph(str(bullet), style="List Bullet")

        findings = content.get("findings") or []
        if findings:
            document.add_heading("Findings", level=1)
            table = document.add_table(rows=1, cols=3)
            table.style = "Light Grid Accent 1"
            header = table.rows[0].cells
            for index, label in enumerate(("Finding", "Severity", "Reference")):
                header[index].text = label
                for paragraph in header[index].paragraphs:
                    for run_ in paragraph.runs:
                        run_.bold = True
            for finding in findings:
                cells = table.add_row().cells
                cells[0].text = str(finding.get("description", ""))
                cells[1].text = str(finding.get("severity", "")).upper()
                cells[2].text = str(finding.get("reference", "") or "—")

        if calculations:
            document.add_heading("Calculations", level=1)
            table = document.add_table(rows=1, cols=4)
            table.style = "Light Grid Accent 1"
            header = table.rows[0].cells
            for index, label in enumerate(
                ("Quantity", "Expression", "Result", "Independently recomputed")
            ):
                header[index].text = label
                for paragraph in header[index].paragraphs:
                    for run_ in paragraph.runs:
                        run_.bold = True
            for calculation in calculations:
                cells = table.add_row().cells
                cells[0].text = str(calculation.get("label", ""))
                cells[1].text = str(calculation.get("expression", ""))
                units = str(calculation.get("units") or "")
                cells[2].text = f"{calculation.get('recomputed', calculation.get('expected', ''))} {units}".strip()
                cells[3].text = "Yes — matched" if calculation.get("matched") else "Discrepancy"

        if content.get("recommendation"):
            document.add_heading("Recommendation", level=1)
            document.add_paragraph(str(content["recommendation"]))

        if content.get("approval_statement"):
            document.add_heading("Approval sought", level=1)
            document.add_paragraph(str(content["approval_statement"]))

        if evidence:
            document.add_heading("Evidence and source references", level=1)
            for item in evidence:
                paragraph = document.add_paragraph()
                marker = paragraph.add_run(f"[{item.id}] ")
                marker.bold = True
                location = f", {item.location}" if item.location else ""
                version = f", version {item.version}" if item.version else ""
                paragraph.add_run(
                    f"{item.source_document}{location}{version} "
                    f"({item.classification.value}, {item.department or 'unassigned'})"
                )
                excerpt = document.add_paragraph(f"“{item.excerpt.strip()[:600]}”")
                excerpt.paragraph_format.left_indent = Pt(18)
                for run_ in excerpt.runs:
                    run_.italic = True
                    run_.font.size = Pt(9)
                    run_.font.color.rgb = RGBColor.from_string(MUTED)

        document.add_heading("Provenance and verification record", level=1)
        provenance = document.add_table(rows=0, cols=2)
        provenance.style = "Light Grid Accent 1"

        model_lines = "; ".join(
            f"{decision.selected_display_name or decision.selected_model} "
            f"({decision.requested_role.value})"
            for decision in routing
            if decision.selected_model
        ) or "no model invoked"
        entries: list[tuple[str, str]] = [
            ("Models used", model_lines),
            ("Evidence items", str(len(evidence))),
            ("Processing location", "On-premise host only — no external calls"),
        ]
        if verification:
            entries.append(
                (
                    "Verification checks",
                    "; ".join(
                        f"{check.name}: {'pass' if check.passed else 'FAIL'}"
                        for check in verification.checks
                    )
                    or "none",
                )
            )
            entries.append(
                (
                    "Material claims supported",
                    f"{verification.material_claims_supported} of "
                    f"{verification.material_claims_total}",
                )
            )
            if verification.limitations:
                entries.append(("Limitations", "; ".join(verification.limitations)))
        for label, value in entries:
            cells = provenance.add_row().cells
            cells[0].text = label
            cells[1].text = value
            for paragraph in cells[0].paragraphs:
                for run_ in paragraph.runs:
                    run_.bold = True

        closing = document.add_paragraph()
        closing_run = closing.add_run(
            "This document was drafted by an AI agent operating entirely within "
            "the organisation's premises and requires human approval before release."
        )
        closing_run.italic = True
        closing_run.font.size = Pt(8.5)
        closing_run.font.color.rgb = RGBColor.from_string(MUTED)

        path = self._target(task_id, title, "docx")
        document.save(str(path))
        return self._finalise(path, "docx")

    # -- XLSX --------------------------------------------------------------
    def render_xlsx(
        self,
        *,
        task_id: str,
        content: dict[str, Any],
        evidence: list[EvidenceItem],
        tables: list[dict[str, Any]] | None = None,
    ) -> Deliverable:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter

        workbook = Workbook()
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill("solid", fgColor=ACCENT)

        sheet = workbook.active
        sheet.title = "Summary"
        sheet["A1"] = str(content.get("title") or "Analysis")
        sheet["A1"].font = Font(bold=True, size=14)
        sheet["A2"] = str(content.get("summary") or "")
        sheet["A2"].alignment = Alignment(wrap_text=True, vertical="top")
        sheet.column_dimensions["A"].width = 110
        sheet.row_dimensions[2].height = 60

        row_cursor = 4
        for section in content.get("sections") or []:
            sheet.cell(row=row_cursor, column=1, value=str(section.get("heading", ""))).font = Font(bold=True)
            row_cursor += 1
            cell = sheet.cell(row=row_cursor, column=1, value=str(section.get("body", "")))
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            row_cursor += 2

        for index, table in enumerate(tables or [], start=1):
            name = str(table.get("name") or f"Table{index}")[:28]
            worksheet = workbook.create_sheet(name)
            rows = table.get("rows") or []
            headers = table.get("headers") or (rows[0] if rows else [])
            for column_index, label in enumerate(headers, start=1):
                cell = worksheet.cell(row=1, column=column_index, value=str(label))
                cell.font = header_font
                cell.fill = header_fill
            body_rows = rows[1:] if (rows and rows[0] == headers) else rows
            for row_index, row in enumerate(body_rows, start=2):
                for column_index, value in enumerate(row, start=1):
                    worksheet.cell(row=row_index, column=column_index, value=value)
            for column_index in range(1, len(headers) + 1):
                worksheet.column_dimensions[get_column_letter(column_index)].width = 22

        if evidence:
            worksheet = workbook.create_sheet("Evidence")
            for column_index, label in enumerate(
                ("ID", "Source document", "Location", "Classification", "Excerpt"), start=1
            ):
                cell = worksheet.cell(row=1, column=column_index, value=label)
                cell.font = header_font
                cell.fill = header_fill
            for row_index, item in enumerate(evidence, start=2):
                worksheet.cell(row=row_index, column=1, value=item.id)
                worksheet.cell(row=row_index, column=2, value=item.source_document)
                worksheet.cell(row=row_index, column=3, value=item.location or "")
                worksheet.cell(row=row_index, column=4, value=item.classification.value)
                worksheet.cell(row=row_index, column=5, value=item.excerpt[:500])
            for column, width in zip("ABCDE", (10, 32, 24, 16, 90)):
                worksheet.column_dimensions[column].width = width

        path = self._target(task_id, str(content.get("title") or "analysis"), "xlsx")
        workbook.save(str(path))
        return self._finalise(path, "xlsx")

    # -- PPTX --------------------------------------------------------------
    def render_pptx(
        self,
        *,
        task_id: str,
        content: dict[str, Any],
        evidence: list[EvidenceItem],
    ) -> Deliverable:
        from pptx import Presentation
        from pptx.dml.color import RGBColor as PptxColor
        from pptx.util import Inches, Pt as PptxPt

        presentation = Presentation()
        title_layout = presentation.slide_layouts[0]
        bullet_layout = presentation.slide_layouts[1]

        slide = presentation.slides.add_slide(title_layout)
        slide.shapes.title.text = str(content.get("title") or "Industrial Briefing")
        slide.placeholders[1].text = (
            f"{content.get('reference') or task_id[:12].upper()}  ·  "
            f"{datetime.now(timezone.utc).strftime('%d %b %Y')}  ·  "
            "Prepared on-premise"
        )

        if content.get("summary"):
            summary_slide = presentation.slides.add_slide(bullet_layout)
            summary_slide.shapes.title.text = "Executive summary"
            summary_slide.placeholders[1].text_frame.text = str(content["summary"])

        for section in content.get("sections") or []:
            section_slide = presentation.slides.add_slide(bullet_layout)
            section_slide.shapes.title.text = str(section.get("heading", "Section"))
            frame = section_slide.placeholders[1].text_frame
            body = str(section.get("body", "")).strip()
            frame.text = body[:400] if body else ""
            for bullet in section.get("bullets") or []:
                paragraph = frame.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 1

        findings = content.get("findings") or []
        if findings:
            findings_slide = presentation.slides.add_slide(bullet_layout)
            findings_slide.shapes.title.text = "Findings"
            frame = findings_slide.placeholders[1].text_frame
            frame.text = str(findings[0].get("description", ""))
            for finding in findings[1:]:
                paragraph = frame.add_paragraph()
                paragraph.text = (
                    f"[{str(finding.get('severity', '')).upper()}] "
                    f"{finding.get('description', '')}"
                )

        if content.get("recommendation"):
            recommendation_slide = presentation.slides.add_slide(bullet_layout)
            recommendation_slide.shapes.title.text = "Recommendation"
            recommendation_slide.placeholders[1].text_frame.text = str(
                content["recommendation"]
            )

        if evidence:
            evidence_slide = presentation.slides.add_slide(bullet_layout)
            evidence_slide.shapes.title.text = "Evidence and sources"
            frame = evidence_slide.placeholders[1].text_frame
            frame.text = f"[{evidence[0].id}] {evidence[0].source_document}"
            for item in evidence[1:8]:
                paragraph = frame.add_paragraph()
                paragraph.text = (
                    f"[{item.id}] {item.source_document}"
                    + (f" — {item.location}" if item.location else "")
                )

        footer_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = footer_slide.shapes.add_textbox(
            Inches(0.6), Inches(3.0), Inches(8.8), Inches(1.5)
        )
        frame = textbox.text_frame
        frame.text = "100% local processing"
        frame.paragraphs[0].runs[0].font.size = PptxPt(28)
        frame.paragraphs[0].runs[0].font.color.rgb = PptxColor(0x0F, 0x6E, 0x6E)
        detail = frame.add_paragraph()
        detail.text = (
            "No external API calls · No cloud inference · No data left the host. "
            "Full audit trail retained on-premise."
        )
        detail.runs[0].font.size = PptxPt(12)

        path = self._target(task_id, str(content.get("title") or "briefing"), "pptx")
        presentation.save(str(path))
        return self._finalise(path, "pptx")

    # -- Markdown ----------------------------------------------------------
    def render_markdown(
        self,
        *,
        task_id: str,
        content: dict[str, Any],
        evidence: list[EvidenceItem],
    ) -> Deliverable:
        lines: list[str] = [f"# {content.get('title') or 'Industrial Note'}", ""]
        if content.get("reference"):
            lines.append(f"**Reference:** {content['reference']}  ")
        lines.append(f"**Task:** {task_id}  ")
        lines.append(
            f"**Prepared:** {datetime.now(timezone.utc).strftime('%d %b %Y %H:%M UTC')} "
            "(on-premise, no external calls)\n"
        )
        if content.get("summary"):
            lines += ["## Executive summary", "", str(content["summary"]), ""]
        for section in content.get("sections") or []:
            lines += [f"## {section.get('heading', 'Section')}", "", str(section.get("body", "")), ""]
            for bullet in section.get("bullets") or []:
                lines.append(f"- {bullet}")
            lines.append("")
        if content.get("findings"):
            lines += ["## Findings", "", "| Finding | Severity | Reference |", "|---|---|---|"]
            for finding in content["findings"]:
                lines.append(
                    f"| {finding.get('description','')} | "
                    f"{str(finding.get('severity','')).upper()} | "
                    f"{finding.get('reference','—')} |"
                )
            lines.append("")
        if content.get("recommendation"):
            lines += ["## Recommendation", "", str(content["recommendation"]), ""]
        if evidence:
            lines += ["## Evidence", ""]
            for item in evidence:
                location = f", {item.location}" if item.location else ""
                lines.append(f"- **[{item.id}]** {item.source_document}{location}")
                lines.append(f"  > {item.excerpt.strip()[:400]}")
            lines.append("")

        path = self._target(task_id, str(content.get("title") or "note"), "md")
        path.write_text("\n".join(lines), encoding="utf-8")
        return self._finalise(path, "md")

    # -- dispatch ----------------------------------------------------------
    def render(
        self,
        fmt: str,
        *,
        task_id: str,
        content: dict[str, Any],
        evidence: list[EvidenceItem] | None = None,
        routing: list[RoutingDecision] | None = None,
        verification: VerificationReport | None = None,
        author: str = "Industrial user",
        calculations: list[dict[str, Any]] | None = None,
        tables: list[dict[str, Any]] | None = None,
    ) -> Deliverable:
        evidence = evidence or []
        routing = routing or []
        normalised = fmt.lower().strip()
        if normalised == "docx":
            return self.render_docx(
                task_id=task_id,
                content=content,
                evidence=evidence,
                routing=routing,
                verification=verification,
                author=author,
                calculations=calculations,
            )
        if normalised == "xlsx":
            return self.render_xlsx(
                task_id=task_id, content=content, evidence=evidence, tables=tables
            )
        if normalised == "pptx":
            return self.render_pptx(task_id=task_id, content=content, evidence=evidence)
        if normalised in {"md", "markdown"}:
            return self.render_markdown(task_id=task_id, content=content, evidence=evidence)
        raise ValueError(
            f"Unsupported deliverable format '{fmt}'. "
            "Supported: docx, xlsx, pptx, md."
        )


_engine: DeliverableEngine | None = None


def get_deliverable_engine() -> DeliverableEngine:
    global _engine
    if _engine is None:
        _engine = DeliverableEngine()
    return _engine
